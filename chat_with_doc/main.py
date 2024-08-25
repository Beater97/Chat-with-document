import ollama 
import chromadb
import pandas as pd
import docx2txt
from pptx import Presentation
import markdown
import os
import PyPDF4
import xml.etree.ElementTree as ET
import json
import yaml
from datetime import datetime
import time

class FileReader:

    def read_excel(self, file_path):
        """
        Legge un file Excel e restituisce una stringa CSV con il contenuto del dataframe.
        """
        try:
            df = pd.read_excel(file_path)
            csv_data = df.to_csv(index=False)  # Converte il dataframe in CSV senza includere l'indice
            return csv_data
        except Exception as e:
            return f"Errore nella lettura del file Excel: {str(e)}"

    def read_txt(self, file_path):
        """
        Legge un file di testo e restituisce il contenuto come stringa.
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            return content
        except Exception as e:
            return f"Errore nella lettura del file TXT: {str(e)}"

    def read_docx(self, file_path):
        """
        Legge un file DOCX e restituisce il contenuto come stringa.
        """
        try:
            content = docx2txt.process(file_path)
            return content
        except Exception as e:
            return f"Errore nella lettura del file DOCX: {str(e)}"

    def read_pptx(self, file_path):
        """
        Legge un file PPTX e restituisce una stringa con il testo di ciascuna diapositiva.
        """
        try:
            presentation = Presentation(file_path)
            slides_text = []
            for slide in presentation.slides:
                slide_text = ''
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + '\n'
                slides_text.append(slide_text)
            return "\n".join(slides_text)
        except Exception as e:
            return f"Errore nella lettura del file PPTX: {str(e)}"

    def read_md(self, file_path):
        """
        Legge un file Markdown e restituisce il contenuto come stringa HTML.
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            html = markdown.markdown(text)
            return html
        except Exception as e:
            return f"Errore nella lettura del file MD: {str(e)}"

    def read_pdf(self, file_path):
        """
        Legge un file PDF e restituisce il testo come stringa.
        """
        try:
            pdf_text = ""
            with open(file_path, 'rb') as file:
                reader = PyPDF4.PdfFileReader(file)
                for page_num in range(reader.numPages):
                    page = reader.getPage(page_num)
                    pdf_text += page.extractText()
            return pdf_text
        except Exception as e:
            return f"Errore nella lettura del file PDF: {str(e)}"

    def read_xml(self, file_path):
        """
        Legge un file XML e restituisce una stringa con il contenuto XML.
        """
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            return ET.tostring(root, encoding='utf-8').decode('utf-8')
        except Exception as e:
            return f"Errore nella lettura del file XML: {str(e)}"

    def read_swagger(self, file_path):
        """
        Legge un file Swagger (in formato YAML o JSON) e restituisce una stringa JSON.
        """
        try:
            extension = os.path.splitext(file_path)[1].lower()
            with open(file_path, 'r', encoding='utf-8') as file:
                if extension == '.yaml' or extension == '.yml':
                    data = yaml.safe_load(file)
                elif extension == '.json':
                    data = json.load(file)
                else:
                    return "Formato non riconosciuto per Swagger: deve essere YAML o JSON."
            return json.dumps(data, indent=4)
        except Exception as e:
            return f"Errore nella lettura del file Swagger: {str(e)}"
        
    def read_file(self, file_path):
        """
        Legge un file in base alla sua estensione e restituisce il contenuto come stringa.
        """
        extension = os.path.splitext(file_path)[1].lower().replace("'","")
        file_path = file_path.replace("'","").replace("& ","")

        if extension == '.xlsx' or extension == '.xls':
            return self.read_excel(file_path)
        elif extension == '.txt':
            return self.read_txt(file_path)
        elif extension == '.docx':
            return self.read_docx(file_path)
        elif extension == '.pptx':
            return self.read_pptx(file_path)
        elif extension == '.md':
            return self.read_md(file_path)
        elif extension == '.pdf':
            return self.read_pdf(file_path)
        elif extension == '.xml':
            return self.read_xml(file_path)
        elif extension == '.json' or extension == '.yaml' or extension == '.yml':
            return self.read_swagger(file_path)
        else:
            return f"Tipo di file non supportato: {extension}"
class DocumentManager:
    def __init__(self, collection_name="docs"):
        self.client = chromadb.Client()
        self.collection = self.client.create_collection(name=collection_name)

    def insert_document(self, document_text, ollama_model="mxbai-embed-large"):
        # Ottieni l'embedding del documento
        response = ollama.embeddings(model=ollama_model, prompt=document_text)
        embedding = response["embedding"]
        
        # Aggiungi il documento alla collezione
        self.collection.add(
            ids=[datetime.now().strftime("%Y%m%d%H%M%S")],
            embeddings=[embedding],
            documents=[document_text]
        )
        print("Documento inserito e analizzato.")

    def retrieve_document(self, query, ollama_model="mxbai-embed-large", n_results=1):
        # Ottieni l'embedding della query
        response = ollama.embeddings(prompt=query, model=ollama_model)
        query_embedding = response["embedding"]
        
        # Esegui la query nella collezione
        results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=n_results
        )
        print(results)
        
        if results and 'documents' in results:
            return results['documents'][0][0]
        else:
            print("Nessun risultato trovato.")
            return None

    def generate_response(self, data, prompt, ollama_model="gwen:0.5b"):

        stream = ollama.chat(
            model=ollama_model,
            messages=[{'role': 'user', 'content': f"Usando questi dati: {data}. Rispondi a questo prompt in italiano: {prompt}"}],
            stream=True
        )
        for chunk in stream:
            print(chunk['message']['content'], end='', flush=True)

reader = FileReader()

doc_manager = DocumentManager()

while True :
    print("\nMenu:")
    print("1. Inserisci documento")
    print("2. Fai una domanda")
    print("3. Esci\n")
    action = input()
    if int(action) == 1 :
        documento = input('dragga e droppa il tuo documento qui e poi clicca invio : ')
        doc_manager.insert_document(reader.read_file(documento))
    if int(action) == 2 :
        prompt =  input('fai la tua domanda : ')
        data = doc_manager.retrieve_document(prompt)
        doc_manager.generate_response(data,prompt)